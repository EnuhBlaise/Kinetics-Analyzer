#!/usr/bin/env python3
"""
Generate a conference-style PowerPoint presentation showcasing the
modular software architecture for kinetic parameter estimation.

Usage:
    python scripts/make_presentation.py

Output:
    output/presentation.pptx
"""

import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── colour palette ──────────────────────────────────────────────
NAVY       = RGBColor(0x1B, 0x2A, 0x4A)
DARK_NAVY  = RGBColor(0x0F, 0x1A, 0x33)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
DARK_GRAY  = RGBColor(0x33, 0x33, 0x33)
TEAL       = RGBColor(0x00, 0x77, 0xB6)
ACCENT_RED = RGBColor(0xEE, 0x66, 0x77)
ACCENT_GRN = RGBColor(0x22, 0x88, 0x33)
MID_BLUE   = RGBColor(0x44, 0x77, 0xAA)
SOFT_BG    = RGBColor(0xE8, 0xEE, 0xF4)

# ── font defaults ───────────────────────────────────────────────
FONT_TITLE   = "Calibri"
FONT_BODY    = "Calibri"
PT_TITLE     = Pt(28)
PT_SUBTITLE  = Pt(18)
PT_BODY      = Pt(16)
PT_SMALL     = Pt(14)
PT_TINY      = Pt(12)

# ── slide dimensions (standard 16:9) ───────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ═══════════════════════════════════════════════════════════════
# Helper utilities
# ═══════════════════════════════════════════════════════════════

def _set_slide_bg(slide, color):
    """Fill the slide background with a solid colour."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, left, top, width, height):
    """Add a textbox and return its text_frame."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    return txBox.text_frame


def _set_font(run, size=PT_BODY, color=DARK_GRAY, bold=False, italic=False,
              name=FONT_BODY):
    """Apply font settings to a run."""
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = name


def _add_paragraph(tf, text, size=PT_BODY, color=DARK_GRAY, bold=False,
                   italic=False, alignment=PP_ALIGN.LEFT, space_after=Pt(6),
                   space_before=Pt(0), level=0):
    """Add a styled paragraph to a text_frame."""
    p = tf.add_paragraph()
    p.alignment = alignment
    p.space_after = space_after
    p.space_before = space_before
    p.level = level
    run = p.add_run()
    run.text = text
    _set_font(run, size=size, color=color, bold=bold, italic=italic)
    return p


def _add_bullet(tf, text, size=PT_BODY, color=DARK_GRAY, level=0,
                bold=False, space_after=Pt(4)):
    """Add a bullet-point paragraph."""
    return _add_paragraph(tf, text, size=size, color=color, bold=bold,
                          level=level, space_after=space_after)


def _add_title_bar(slide, title_text, subtitle_text=None):
    """Add a navy title bar across the top of a content slide."""
    # Title band
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(1.15)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = title_text
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    run = tf.paragraphs[0].runs[0]
    _set_font(run, size=PT_TITLE, color=WHITE, bold=True)

    # left padding via indentation
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.15)

    if subtitle_text:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = subtitle_text
        _set_font(r, size=PT_SMALL, color=RGBColor(0xBB, 0xCC, 0xDD),
                  italic=True)


def _content_frame(slide, left=Inches(0.7), top=Inches(1.4),
                   width=Inches(11.9), height=Inches(5.6)):
    """Return a text_frame positioned in the body area of a content slide."""
    tf = _add_textbox(slide, left, top, width, height)
    tf.word_wrap = True
    return tf


def _add_box(slide, left, top, width, height, fill_color, text,
             font_size=PT_SMALL, font_color=WHITE, bold=True,
             border_color=None):
    """Add a rounded-rectangle box with centred text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)

    # vertical centering
    tf.auto_size = None
    shape.text_frame._txBody.attrib[
        '{http://schemas.openxmlformats.org/drawingml/2006/main}anchor'
    ] = 'ctr' if False else 'ctr'

    run = tf.paragraphs[0].add_run()
    run.text = text
    _set_font(run, size=font_size, color=font_color, bold=bold)
    return shape


def _add_arrow_down(slide, cx, top, length=Inches(0.35)):
    """Draw a small downward arrow at (cx, top)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW, cx - Inches(0.15), top,
        Inches(0.3), length
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = TEAL
    shape.line.fill.background()


# ═══════════════════════════════════════════════════════════════
# Slide builders
# ═══════════════════════════════════════════════════════════════

def slide_01_title(prs):
    """Slide 1 — Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_slide_bg(slide, DARK_NAVY)

    # Main title
    tf = _add_textbox(slide, Inches(1), Inches(2.0), Inches(11.3), Inches(1.5))
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Modular Software Architecture in\nComputational Biology"
    _set_font(run, size=Pt(36), color=WHITE, bold=True)

    # Subtitle
    tf2 = _add_textbox(slide, Inches(1.5), Inches(3.8), Inches(10.3), Inches(1.2))
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = ("A Case Study in Nonlinear Kinetic Parameter Estimation\n"
                 "for Microbial Growth Modeling")
    _set_font(run2, size=PT_SUBTITLE, color=RGBColor(0xAA, 0xCC, 0xEE),
              italic=True)

    # Decorative line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(4), Inches(3.55), Inches(5.3), Pt(2)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()

    # Footer placeholder
    tf3 = _add_textbox(slide, Inches(2), Inches(6.2), Inches(9.3), Inches(0.6))
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run()
    r3.text = "Kinetic Parameter Estimation Software v2.0"
    _set_font(r3, size=PT_SMALL, color=RGBColor(0x88, 0x99, 0xAA))


def slide_02_motivation(prs):
    """Slide 2 — The Problem."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, "The Problem: Monolithic Parameter Estimation")

    tf = _content_frame(slide)
    bullets = [
        "Nonlinear ODE parameter estimation is ubiquitous in biology, "
        "environmental engineering, and chemical kinetics",
        "Most labs use monolithic scripts: model + fitting + I/O tightly coupled",
        "Changing a model requires rewriting the optimization loop",
        "Adding a new statistical metric touches fitting code",
        "Reproducibility and extensibility suffer",
        "No systematic uncertainty quantification or model comparison",
    ]
    for b in bullets:
        _add_bullet(tf, b, size=PT_BODY, color=DARK_GRAY)


def slide_03_approach(prs):
    """Slide 3 — Design Philosophy."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, "Design Philosophy: Separation of Concerns")

    tf = _content_frame(slide)
    bullets = [
        "Core mathematics isolated from I/O, optimization, and workflows",
        "Each layer depends only on the layer below",
        "New models, optimizers, or CI methods plug in without modifying existing code",
        "Configuration-driven: new substrates via JSON, no code changes",
        "Tested: 43+ unit tests covering core, fitting, and workflows",
        "Multiple user interfaces: CLI scripts + Streamlit web app",
    ]
    for b in bullets:
        _add_bullet(tf, b, size=PT_BODY, color=DARK_GRAY)


def slide_04_architecture(prs):
    """Slide 4 — Layered Architecture diagram."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, "Layered Architecture")

    box_w = Inches(8.0)
    box_h = Inches(0.7)
    left_main = Inches(2.6)
    top_start = Inches(1.6)
    gap = Inches(0.95)

    # Layer 1 - top
    _add_box(slide, left_main, top_start, box_w, box_h,
             RGBColor(0x3A, 0x86, 0xA8),
             "scripts/  &  streamlit_app/  \u2014  User Interfaces (CLI + Web)")
    _add_arrow_down(slide, Inches(6.6), top_start + box_h + Pt(2))

    # Layer 2
    y2 = top_start + gap
    _add_box(slide, left_main, y2, box_w, box_h,
             RGBColor(0x00, 0x77, 0xB6),
             "workflows/  \u2014  Orchestration (BaseWorkflow \u2192 Concrete Workflows)")
    _add_arrow_down(slide, Inches(6.6), y2 + box_h + Pt(2))

    # Layer 3
    y3 = top_start + 2 * gap
    _add_box(slide, left_main, y3, box_w, box_h,
             RGBColor(0x00, 0x59, 0x8A),
             "src/fitting/  \u2014  Optimization, Statistics, Objective Functions")
    _add_arrow_down(slide, Inches(6.6), y3 + box_h + Pt(2))

    # Layer 4 - bottom
    y4 = top_start + 3 * gap
    _add_box(slide, left_main, y4, box_w, box_h,
             NAVY,
             "src/core/  \u2014  Monod Kinetics, ODE Systems, RK45 Solver")

    # Side boxes
    side_w = Inches(2.0)
    side_h = Inches(1.3)
    side_left = Inches(0.3)

    _add_box(slide, side_left, top_start + Inches(0.3), side_w, side_h,
             RGBColor(0x6C, 0x95, 0x7F),
             "src/io/\nConfig, Data,\nResults I/O",
             font_size=PT_TINY)

    _add_box(slide, side_left, top_start + Inches(2.0), side_w, side_h,
             RGBColor(0x8E, 0x7C, 0xA8),
             "src/utils/\nPlotting,\nValidation,\nConversions",
             font_size=PT_TINY)

    # Key label
    tf = _add_textbox(slide, Inches(0.5), Inches(6.2), Inches(12), Inches(0.5))
    _add_paragraph(tf, "Each layer depends only on the layer below  \u2014  "
                   "side modules are shared utilities",
                   size=PT_SMALL, color=TEAL, italic=True,
                   alignment=PP_ALIGN.CENTER)


def slide_05_math_core(prs):
    """Slide 5 — Mathematical Core."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, "Core Layer: Kinetic Models")

    tf = _content_frame(slide, top=Inches(1.35))

    _add_paragraph(tf, "Monod (Haldane) equation:", size=PT_BODY,
                   color=NAVY, bold=True, space_after=Pt(2))
    _add_paragraph(tf, "    q = q_max \u00b7 S / (Ks + S + S\u00b2/Ki)",
                   size=PT_BODY, color=DARK_GRAY, space_after=Pt(10))

    _add_paragraph(tf, "Dual Monod (oxygen limitation):", size=PT_BODY,
                   color=NAVY, bold=True, space_after=Pt(2))
    _add_paragraph(tf, "    q_dual = q \u00b7 O\u2082 / (K_O2 + O\u2082)",
                   size=PT_BODY, color=DARK_GRAY, space_after=Pt(10))

    _add_paragraph(tf, "Lag phase factor (sigmoid):", size=PT_BODY,
                   color=NAVY, bold=True, space_after=Pt(2))
    _add_paragraph(tf, "    f(t) = 1 / (1 + exp(\u2013k \u00b7 (t \u2013 \u03bb/2) / \u03bb))",
                   size=PT_BODY, color=DARK_GRAY, space_after=Pt(14))

    _add_paragraph(tf, "Three ODE systems \u2014 all inherit from BaseODE:",
                   size=PT_BODY, color=NAVY, bold=True, space_after=Pt(4))

    models = [
        ("SingleMonodODE", "2 states (S, X), 5 parameters"),
        ("DualMonodODE", "3 states (S, X, O\u2082), 7 parameters"),
        ("DualMonodLagODE", "3 states + lag, 8 parameters"),
    ]
    for name, desc in models:
        p = tf.add_paragraph()
        p.space_after = Pt(3)
        r1 = p.add_run()
        r1.text = f"    {name}: "
        _set_font(r1, size=PT_BODY, color=TEAL, bold=True)
        r2 = p.add_run()
        r2.text = desc
        _set_font(r2, size=PT_BODY, color=DARK_GRAY)


def slide_06_biological(prs, image_path):
    """Slide 6 — Biological Schematic (embed image)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, "Biological Model")

    if os.path.exists(image_path):
        # Centre the image
        slide.shapes.add_picture(
            image_path, Inches(2.5), Inches(1.5), Inches(8.3), Inches(5.5)
        )
    else:
        tf = _content_frame(slide)
        _add_paragraph(tf, f"[Image not found: {image_path}]",
                       size=PT_BODY, color=ACCENT_RED, italic=True)


def slide_07_ode_hierarchy(prs):
    """Slide 7 — ODE Class Hierarchy."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, "Extensible ODE System Design")

    # BaseODE box (top centre)
    base_left = Inches(4.0)
    base_top = Inches(1.6)
    base_w = Inches(5.3)
    base_h = Inches(1.1)
    _add_box(slide, base_left, base_top, base_w, base_h, NAVY,
             "BaseODE (ABC)\nderivatives()  |  state_names  |  state_units  |  n_states",
             font_size=PT_SMALL)

    # Arrows down to children
    children_tops = base_top + base_h + Inches(0.5)
    for cx in [Inches(2.6), Inches(6.0), Inches(9.4)]:
        _add_arrow_down(slide, cx, base_top + base_h + Pt(4), Inches(0.35))

    # Three child boxes
    child_w = Inches(3.4)
    child_h = Inches(1.1)
    child_y = children_tops + Inches(0.1)

    _add_box(slide, Inches(0.9), child_y, child_w, child_h,
             MID_BLUE,
             "SingleMonodODE\n2 states (S, X)\n5 parameters",
             font_size=PT_TINY)

    _add_box(slide, Inches(4.8), child_y, child_w, child_h,
             RGBColor(0x00, 0x59, 0x8A),
             "DualMonodODE\n3 states (S, X, O\u2082)\n7 parameters",
             font_size=PT_TINY)

    _add_box(slide, Inches(8.7), child_y, child_w, child_h,
             RGBColor(0x2A, 0x40, 0x66),
             "DualMonodLagODE\n3 states + lag\n8 parameters",
             font_size=PT_TINY)

    # Key point
    tf = _add_textbox(slide, Inches(1.0), Inches(5.2), Inches(11.3), Inches(1.5))
    tf.word_wrap = True
    _add_paragraph(
        tf,
        "Adding a new model = one new class, zero changes elsewhere",
        size=PT_SUBTITLE, color=TEAL, bold=True,
        alignment=PP_ALIGN.CENTER
    )
    bullets = [
        "Implement derivatives(), state_names, state_units, n_states",
        "Automatically works with all fitting, statistics, and visualization code",
    ]
    for b in bullets:
        _add_bullet(tf, b, size=PT_SMALL, color=DARK_GRAY, level=0)


def slide_08_fitting(prs):
    """Slide 8 — Fitting Layer."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, "Fitting Layer: Objective Functions & Optimization")

    tf = _content_frame(slide)

    items = [
        ("ObjectiveFunction",
         "Wraps ODE solver, computes normalized SSE per condition"),
        ("GlobalObjectiveFunction",
         "Sums individual losses:  J = \u03a3 L\u1d62  across all conditions"),
        ("ParameterOptimizer",
         "Pluggable algorithms \u2014 L-BFGS-B (gradient-based) and "
         "Differential Evolution (global, gradient-free)"),
    ]
    for label, desc in items:
        p = tf.add_paragraph()
        p.space_after = Pt(6)
        r1 = p.add_run()
        r1.text = f"{label}:  "
        _set_font(r1, size=PT_BODY, color=TEAL, bold=True)
        r2 = p.add_run()
        r2.text = desc
        _set_font(r2, size=PT_BODY, color=DARK_GRAY)

    _add_paragraph(tf, "", size=Pt(6))  # spacer
    _add_paragraph(tf, "Normalized errors for scale-invariant fitting:",
                   size=PT_BODY, color=NAVY, bold=True, space_after=Pt(2))
    _add_paragraph(
        tf,
        "    SSE_norm = SSE_S / range(S)\u00b2  +  SSE_X / range(X)\u00b2",
        size=PT_BODY, color=DARK_GRAY, space_after=Pt(8)
    )
    _add_paragraph(
        tf,
        "Substrate and biomass statistics computed separately, "
        "then combined via weighted average",
        size=PT_SMALL, color=DARK_GRAY, italic=True
    )


def slide_09_two_stage(prs):
    """Slide 9 — Two-Stage Estimation Strategy."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, "Two-Stage Global Parameter Estimation")

    tf = _content_frame(slide)

    stages = [
        ("Stage 1 \u2014 Individual Fits",
         "Fit each concentration condition independently \u2192 local optima per condition"),
        ("Stage 2 \u2014 Global Optimization",
         "Median of individual fits \u2192 initial guess for global minimization of "
         "J = \u03a3\u1d62 L\u1d62 via L-BFGS-B"),
    ]
    for title, desc in stages:
        _add_paragraph(tf, title, size=PT_BODY, color=NAVY, bold=True,
                       space_after=Pt(2))
        _add_paragraph(tf, desc, size=PT_BODY, color=DARK_GRAY,
                       space_after=Pt(14))

    _add_paragraph(tf, "Advantages:", size=PT_BODY, color=NAVY, bold=True,
                   space_after=Pt(4))
    advantages = [
        "Robust to outlier conditions",
        "Provides per-condition diagnostics and confidence intervals",
        "Individual fits serve as quality check before global estimation",
        "Residual analysis (normality, autocorrelation) per condition",
    ]
    for a in advantages:
        _add_bullet(tf, a, size=PT_BODY, color=DARK_GRAY)


def slide_10_confidence(prs):
    """Slide 10 — Pluggable Uncertainty Quantification."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, "Pluggable Uncertainty Quantification")

    tf = _content_frame(slide)
    _add_paragraph(tf, "Three CI methods (strategy pattern, selected via --ci-method):",
                   size=PT_BODY, color=NAVY, bold=True, space_after=Pt(10))

    methods = [
        ("Hessian (classical)",
         "Invert Fisher information at optimum \u2192 fast, assumes normality. "
         "Covariance = (H/2)\u207b\u00b9 \u00b7 s\u00b2"),
        ("Hessian (log-space)",
         "Mixed log/linear transform for positive parameters \u2192 asymmetric CIs, "
         "better for skewed distributions"),
        ("MCMC (Metropolis)",
         "Random-walk Metropolis with multi-chain diagnostics (R\u0302, ESS) \u2192 "
         "full posterior, no distributional assumptions"),
    ]
    for title, desc in methods:
        p = tf.add_paragraph()
        p.space_after = Pt(10)
        r1 = p.add_run()
        r1.text = f"{title}:  "
        _set_font(r1, size=PT_BODY, color=TEAL, bold=True)
        r2 = p.add_run()
        r2.text = desc
        _set_font(r2, size=PT_BODY, color=DARK_GRAY)

    _add_paragraph(tf, "", size=Pt(4))
    _add_paragraph(
        tf,
        "Adding a new CI method = one function, register in dispatcher",
        size=PT_SMALL, color=TEAL, italic=True, bold=True,
        alignment=PP_ALIGN.CENTER
    )


def slide_11_workflow(prs):
    """Slide 11 — Workflow Orchestration."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, "Workflow Layer: Composable Pipelines")

    tf = _content_frame(slide)

    _add_paragraph(tf, "BaseWorkflow (ABC) defines the pipeline:",
                   size=PT_BODY, color=NAVY, bold=True, space_after=Pt(4))

    steps = [
        "_create_ode_system()  \u2192  _setup_objective()  \u2192  "
        "_run_optimization()  \u2192  _generate_report()"
    ]
    for s in steps:
        _add_paragraph(tf, f"    {s}", size=PT_SMALL, color=TEAL,
                       space_after=Pt(12))

    _add_paragraph(tf, "Concrete workflows:", size=PT_BODY, color=NAVY,
                   bold=True, space_after=Pt(4))

    workflows = [
        ("SingleMonodWorkflow", "5 params, no O\u2082 dynamics"),
        ("DualMonodWorkflow", "7 params, with O\u2082 and reaeration"),
        ("DualMonodLagWorkflow", "8 params, with lag phase"),
        ("IndividualConditionWorkflow",
         "Fits each condition, computes global params, CIs, diagnostics"),
    ]
    for name, desc in workflows:
        p = tf.add_paragraph()
        p.space_after = Pt(4)
        r1 = p.add_run()
        r1.text = f"    {name}:  "
        _set_font(r1, size=PT_BODY, color=TEAL, bold=True)
        r2 = p.add_run()
        r2.text = desc
        _set_font(r2, size=PT_BODY, color=DARK_GRAY)

    _add_paragraph(tf, "", size=Pt(6))
    _add_paragraph(
        tf,
        "Adding a new model workflow = subclass + register in factory map",
        size=PT_SMALL, color=TEAL, italic=True, bold=True,
        alignment=PP_ALIGN.CENTER
    )


def slide_12_model_selection(prs):
    """Slide 12 — Model Selection & Comparison."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, "Systematic Model Comparison")

    tf = _content_frame(slide)

    _add_paragraph(tf, "AIC-based model selection across candidate models:",
                   size=PT_BODY, color=NAVY, bold=True, space_after=Pt(8))

    formulas = [
        "AIC = n \u00b7 ln(SSE/n) + 2k     (with small-sample correction AICc)",
        "BIC = n \u00b7 ln(SSE/n) + k \u00b7 ln(n)",
        "Akaike weights:  w\u1d62 = exp(\u20130.5\u00b7\u0394\u1d62) / \u03a3 exp(\u20130.5\u00b7\u0394\u2c7c)",
    ]
    for f in formulas:
        _add_paragraph(tf, f"    {f}", size=PT_BODY, color=DARK_GRAY,
                       space_after=Pt(6))

    _add_paragraph(tf, "", size=Pt(6))
    _add_paragraph(tf, "Automated comparison pipeline:",
                   size=PT_BODY, color=NAVY, bold=True, space_after=Pt(4))

    bullets = [
        "Run all candidate models (single, dual, dual+lag) per substrate",
        "Rank by AIC, BIC, R\u00b2, NRMSE",
        "Compute Akaike weights (probability of being the best model)",
        "Master results table aggregates across substrates",
    ]
    for b in bullets:
        _add_bullet(tf, b, size=PT_BODY, color=DARK_GRAY)


def slide_13_config(prs):
    """Slide 13 — Configuration-Driven Extensibility."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, "Adding a New Substrate: Zero Code Changes")

    tf = _content_frame(slide, width=Inches(5.5))

    _add_paragraph(tf, "JSON config file specifies everything:",
                   size=PT_BODY, color=NAVY, bold=True, space_after=Pt(6))

    bullets = [
        "Parameter bounds and initial guesses",
        "Oxygen settings (saturation, reaeration rate)",
        "Simulation time and resolution",
        "Substrate metadata (name, MW, units)",
    ]
    for b in bullets:
        _add_bullet(tf, b, size=PT_BODY, color=DARK_GRAY)

    _add_paragraph(tf, "", size=Pt(6))
    _add_paragraph(tf, "Same pipeline handles glucose, xylose,\n"
                   "syringic acid, p-hydroxybenzoic acid, etc.",
                   size=PT_BODY, color=DARK_GRAY, space_after=Pt(10))
    _add_paragraph(
        tf,
        "New substrate = one JSON file, zero code changes",
        size=PT_SMALL, color=TEAL, italic=True, bold=True,
    )

    # JSON snippet box on the right
    json_text = (
        '{\n'
        '  "substrate": {\n'
        '    "name": "Glucose",\n'
        '    "molecular_weight": 180.16\n'
        '  },\n'
        '  "initial_guesses": {\n'
        '    "qmax": 20.0,\n'
        '    "Ks": 400.0, "Ki": 10000.0,\n'
        '    "Y": 0.4, "b_decay": 0.0005\n'
        '  },\n'
        '  "bounds": {\n'
        '    "qmax": [15.0, 25.0],\n'
        '    "Ks": [350.0, 450.0],\n'
        '    ...\n'
        '  }\n'
        '}'
    )
    json_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.2)
    )
    json_box.fill.solid()
    json_box.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)
    json_box.line.color.rgb = TEAL
    json_box.line.width = Pt(1.5)

    jtf = json_box.text_frame
    jtf.word_wrap = True
    jtf.margin_left = Inches(0.2)
    jtf.margin_top = Inches(0.15)

    # Title in box
    p_title = jtf.paragraphs[0]
    r_t = p_title.add_run()
    r_t.text = "config/substrates/glucose.json"
    _set_font(r_t, size=PT_TINY, color=TEAL, bold=True)

    for line in json_text.split('\n'):
        p = jtf.add_paragraph()
        p.space_after = Pt(0)
        p.space_before = Pt(0)
        r = p.add_run()
        r.text = line
        _set_font(r, size=Pt(11), color=DARK_GRAY, name="Courier New")


def slide_14_results(prs):
    """Slide 14 — Results & Diagnostics."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, "Automated Diagnostic Output")

    tf = _content_frame(slide)

    _add_paragraph(tf, "Per-condition diagnostic suite:",
                   size=PT_BODY, color=NAVY, bold=True, space_after=Pt(6))

    bullets = [
        "Individual fits with experimental data overlay",
        "Residual diagnostics (normality, autocorrelation, skewness, kurtosis)",
        "Parameter comparison across conditions (box plots, median \u00b1 IQR)",
        "Confidence interval visualization (Hessian, log-Hessian, or MCMC)",
        "Goodness-of-fit summary: R\u00b2, RMSE, NRMSE (separate for substrate/biomass)",
    ]
    for b in bullets:
        _add_bullet(tf, b, size=PT_BODY, color=DARK_GRAY)

    _add_paragraph(tf, "", size=Pt(8))
    _add_paragraph(tf, "Automated reporting:", size=PT_BODY, color=NAVY,
                   bold=True, space_after=Pt(4))

    reports = [
        "PDF reports generated automatically after each workflow run",
        "JSON results files for programmatic downstream analysis",
        "Publication-quality figures (PNG + SVG at 300 DPI)",
        "Master results table for cross-substrate comparison",
    ]
    for r in reports:
        _add_bullet(tf, r, size=PT_BODY, color=DARK_GRAY)


def slide_15_engineering(prs):
    """Slide 15 — Software Engineering Practices."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, "Quality & Reproducibility")

    tf = _content_frame(slide)

    bullets = [
        "43+ unit tests covering core kinetics, fitting, and workflows (pytest)",
        "Type hints throughout for IDE support and static analysis",
        "Comprehensive logging and performance monitoring (psutil)",
        "Streamlit web interface for non-programmers",
        "Version tracking and automated PDF reporting",
        "Modular I/O: data loaders, config loaders, results writers",
        "Flake8 linting for code style consistency",
    ]
    for b in bullets:
        _add_bullet(tf, b, size=PT_BODY, color=DARK_GRAY)

    _add_paragraph(tf, "", size=Pt(8))

    # Tech stack boxes
    stack_y = Inches(5.4)
    techs = [
        ("NumPy / SciPy", RGBColor(0x3A, 0x86, 0xA8)),
        ("Matplotlib / Seaborn", ACCENT_GRN),
        ("Pandas", RGBColor(0x8E, 0x7C, 0xA8)),
        ("Pytest", ACCENT_RED),
        ("Streamlit", RGBColor(0xFF, 0x4B, 0x4B)),
    ]
    x = Inches(0.7)
    for label, color in techs:
        _add_box(slide, x, stack_y, Inches(2.2), Inches(0.5), color, label,
                 font_size=PT_TINY, font_color=WHITE)
        x += Inches(2.45)


def slide_16_summary(prs):
    """Slide 16 — Summary / Conclusions."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, DARK_NAVY)

    # Title
    tf_title = _add_textbox(slide, Inches(0.7), Inches(0.5),
                            Inches(11.9), Inches(1.0))
    p = tf_title.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "Summary"
    _set_font(r, size=Pt(32), color=WHITE, bold=True)

    # Decorative line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.35), Inches(3.0), Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()

    tf = _add_textbox(slide, Inches(0.7), Inches(1.7), Inches(11.9), Inches(4.5))
    tf.word_wrap = True

    bullets = [
        "Deliberate separation of concerns \u2192 extensible, testable, reproducible",
        "New models:  one ODE class + one workflow subclass",
        "New substrates:  one JSON file, zero code changes",
        "New CI methods:  one strategy function, register in dispatcher",
        "Two-stage estimation with per-condition diagnostics for robust results",
        "Framework applicable beyond microbial kinetics to any ODE-based "
        "parameter estimation problem",
    ]
    for b in bullets:
        _add_bullet(tf, b, size=PT_SUBTITLE, color=WHITE, space_after=Pt(10))

    # Footer
    tf_foot = _add_textbox(slide, Inches(0.7), Inches(6.4),
                           Inches(11.9), Inches(0.6))
    p2 = tf_foot.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "Kinetic Parameter Estimation Software v2.0  |  github.com/your-repo"
    _set_font(r2, size=PT_SMALL, color=RGBColor(0x88, 0x99, 0xAA))


# ═══════════════════════════════════════════════════════════════
# Main
# ═══════════════════════════════════════════════════════════════

def main():
    # Resolve paths relative to project root
    project_root = Path(__file__).resolve().parent.parent
    image_path = str(project_root / "docs" / "image.png")
    output_dir = project_root / "output"
    output_dir.mkdir(parents=True, exist_ok=True)
    output_path = str(output_dir / "presentation.pptx")

    prs = Presentation()
    # Set 16:9 widescreen
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    print("Building presentation...")

    slide_01_title(prs)
    print("  [1/16] Title slide")

    slide_02_motivation(prs)
    print("  [2/16] Motivation")

    slide_03_approach(prs)
    print("  [3/16] Design philosophy")

    slide_04_architecture(prs)
    print("  [4/16] Architecture diagram")

    slide_05_math_core(prs)
    print("  [5/16] Mathematical core")

    slide_06_biological(prs, image_path)
    print("  [6/16] Biological schematic")

    slide_07_ode_hierarchy(prs)
    print("  [7/16] ODE class hierarchy")

    slide_08_fitting(prs)
    print("  [8/16] Fitting layer")

    slide_09_two_stage(prs)
    print("  [9/16] Two-stage estimation")

    slide_10_confidence(prs)
    print("  [10/16] Confidence intervals")

    slide_11_workflow(prs)
    print("  [11/16] Workflow orchestration")

    slide_12_model_selection(prs)
    print("  [12/16] Model comparison")

    slide_13_config(prs)
    print("  [13/16] Configuration extensibility")

    slide_14_results(prs)
    print("  [14/16] Results & diagnostics")

    slide_15_engineering(prs)
    print("  [15/16] Software engineering")

    slide_16_summary(prs)
    print("  [16/16] Summary")

    prs.save(output_path)
    print(f"\nPresentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
